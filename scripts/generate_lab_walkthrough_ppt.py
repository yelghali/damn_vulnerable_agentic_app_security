from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations"
OUT_FILE = OUT_DIR / "zava-agentic-app-security-lab-walkthrough.pptx"
ASSETS = ROOT / "docs" / "assets"
SHOTS = ASSETS / "screenshots"

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(247, 249, 246)
INK = RGBColor(24, 32, 38)
MUTED = RGBColor(86, 98, 105)
DARK = RGBColor(19, 38, 46)
TEAL = RGBColor(0, 132, 122)
CORAL = RGBColor(219, 76, 63)
AMBER = RGBColor(237, 177, 75)
GREEN = RGBColor(59, 142, 95)
BLUE = RGBColor(46, 93, 166)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(230, 241, 238)
LINE = RGBColor(205, 215, 211)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Consolas"


def rgb(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def add_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False, font=BODY_FONT,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06 if margin else 0)
    tf.margin_right = Inches(0.06 if margin else 0)
    tf.margin_top = Inches(0.03 if margin else 0)
    tf.margin_bottom = Inches(0.03 if margin else 0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, kicker=None, dark=False):
    if kicker:
        add_text(slide, kicker.upper(), 0.65, 0.38, 5.2, 0.3, 9, TEAL if not dark else AMBER, bold=True, margin=False)
    add_text(slide, title, 0.62, 0.68, 8.7, 0.7, 27, WHITE if dark else INK, bold=True, font=TITLE_FONT, margin=False)


def add_footer(slide, text="Zava Wealth Advisor security lab"):
    add_text(slide, text, 0.65, 7.08, 7.0, 0.22, 7.5, MUTED, margin=False)


def add_card(slide, x, y, w, h, title, body, accent=TEAL, title_size=14, body_size=10.5):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = LINE
    shp.line.width = Pt(0.8)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.14, w - 0.3, 0.28, title_size, INK, bold=True, margin=False)
    add_text(slide, body, x + 0.18, y + 0.48, w - 0.3, h - 0.58, body_size, MUTED, margin=False)


def add_pill(slide, text, x, y, w, color, text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    add_text(slide, text, x, y + 0.06, w, 0.18, 8.2, text_color, bold=True, align=PP_ALIGN.CENTER, margin=False)


def add_image_contain(slide, path: Path, x, y, w, h):
    if not path.exists():
        add_card(slide, x, y, w, h, "Missing asset", str(path.relative_to(ROOT)), CORAL)
        return
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_flow(slide, labels, x, y, w, colors):
    step_w = w / len(labels)
    for i, label in enumerate(labels):
        sx = x + i * step_w
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(sx), Inches(y), Inches(step_w - 0.16), Inches(0.75))
        shp.fill.solid()
        shp.fill.fore_color.rgb = colors[i % len(colors)]
        shp.line.fill.background()
        add_text(slide, label, sx + 0.08, y + 0.17, step_w - 0.32, 0.3, 11, WHITE, bold=True, align=PP_ALIGN.CENTER, margin=False)
        if i < len(labels) - 1:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(sx + step_w - 0.23), Inches(y + 0.23), Inches(0.28), Inches(0.28))
            line.fill.solid()
            line.fill.fore_color.rgb = MUTED
            line.line.fill.background()


def source_footer(slide, text):
    add_text(slide, "Learn: " + text, 0.65, 6.82, 12.0, 0.2, 7.4, MUTED, margin=False)


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK)
    s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(5.7), W, Inches(1.8)).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = TEAL
    add_text(s, "Hardening a Damn\nVulnerable Agentic AI App", 0.72, 0.72, 7.25, 1.35, 29, WHITE, True, TITLE_FONT, margin=False)
    add_text(s, "Zava Wealth Advisor lab walkthrough", 0.76, 1.85, 7.0, 0.45, 18, RGBColor(196, 230, 223), margin=False)
    add_text(s, "Break V1-V11, then add Azure security layers: Foundry guardrails, PII redaction, secure tools/MCP, Entra + Search ACLs, APIM gateway, governance, evaluations, Purview, red teaming.", 0.78, 2.55, 7.05, 1.28, 15.5, WHITE, margin=False)
    add_image_contain(s, SHOTS / "01-app-overview-vulnerable.png", 8.45, 0.75, 3.95, 3.45)
    add_pill(s, "No fluff: exploit -> control -> proof", 0.78, 5.98, 3.15, AMBER, DARK)
    add_pill(s, "Hosted app currently forced to Foundry", 4.15, 5.98, 3.3, GREEN)

    # 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Lab Objective", "why this lab exists")
    add_card(s, 0.7, 1.6, 3.8, 1.55, "What learners do", "Run a real multi-agent finance app, exploit each missing control, then turn on the Azure layer and verify the exploit is dead.", TEAL, 16, 12)
    add_card(s, 4.75, 1.6, 3.8, 1.55, "What they learn", "Where each control belongs: model, API, identity, retrieval, tool boundary, runtime, data governance, or test gate.", BLUE, 16, 12)
    add_card(s, 8.8, 1.6, 3.8, 1.55, "What success means", "A before/after proof for V1-V11, not a checkbox list. Every mitigation is tied to code, Azure config, and a user-visible result.", GREEN, 16, 12)
    add_flow(s, ["Break", "Fix", "Inspect", "Verify"], 1.15, 4.05, 11.0, [CORAL, TEAL, BLUE, GREEN])
    add_footer(s)

    # 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "The App: Zava Wealth Advisor", "what is being secured")
    add_card(s, 0.7, 1.55, 3.3, 4.6, "Multi-agent finance assistant", "Orchestrator routes to Knowledge/RAG, Accounts, Transactions, and Reporting agents. It handles account balances, credit scores, statements, transfers, and reports.", TEAL, 15, 11.2)
    add_card(s, 4.25, 1.55, 3.3, 4.6, "Sensitive data by design", "Names, customer IDs, account numbers, balances, credit scores, addresses, SSNs, private documents, and state-changing financial actions.", CORAL, 15, 11.2)
    add_card(s, 7.8, 1.55, 4.75, 4.6, "Security boundary map", "Model safety protects prompts and completions. Identity determines customer scope. Search ACLs trim documents. Tool policies gate actions. Gateway/runtime controls protect endpoints. Governance/evals prove drift is caught.", BLUE, 15, 11.2)
    add_footer(s)

    # 4
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Vulnerable Baseline", "part 1")
    add_image_contain(s, SHOTS / "01-app-overview-vulnerable.png", 0.7, 1.45, 5.7, 4.75)
    add_card(s, 6.7, 1.5, 5.7, 1.0, "Baseline route", "Browser -> FastAPI app -> local model/local data. No platform guardrails, no trusted identity, no document trimming, no gateway.", CORAL, 15, 11)
    add_card(s, 6.7, 2.75, 2.7, 1.35, "Learner action", "Click exploit chips for V1-V11 and read event traces under the answers.", AMBER, 14, 10.5)
    add_card(s, 9.7, 2.75, 2.7, 1.35, "Facilitator point", "Do not skip Part 1. The controls only matter after learners see the failure.", TEAL, 14, 10.5)
    add_card(s, 6.7, 4.35, 5.7, 1.55, "Baseline model now explicit", "The hosted shared app is forced to Foundry for cohort stability. Local ACA/Phi remains available as an explicit model-route choice for vulnerable-model demos.", BLUE, 14, 10.5)
    add_footer(s)

    # 5
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "V1-V11 Vulnerability Map", "decoder ring")
    vulns = [
        ("V1", "Ungoverned model", CORAL), ("V2", "No guardrails", CORAL), ("V3", "PII leakage", AMBER),
        ("V4", "Overpermissioned tools", AMBER), ("V5", "Broken customer auth", BLUE), ("V6", "Data poisoning", CORAL),
        ("V7", "Insecure runtime", BLUE), ("V8", "Unsafe code execution", AMBER), ("V9", "Insecure MCP", AMBER),
        ("V10", "No AI gateway", BLUE), ("V11", "A2A poisoning", AMBER),
    ]
    for i, (code, label, color) in enumerate(vulns):
        row, col = divmod(i, 4)
        x, y = 0.75 + col * 3.05, 1.45 + row * 1.3
        add_pill(s, code, x, y, 0.65, color)
        add_text(s, label, x + 0.78, y + 0.04, 2.0, 0.32, 11.5, INK, True, margin=False)
    add_card(s, 0.75, 5.45, 11.9, 0.82, "Use this slide as the transition map", "Modules are grouped by Azure security layer, not by vulnerability number. Module 4 intentionally closes V4, V8, V9, and V11 together because they share the tool/action boundary.", TEAL, 14, 10.5)
    add_footer(s)

    layer_slides = [
        ("M1-M2 · Foundry Guardrails", "Responsible AI, Prompt Shields, groundedness", "Exploit", "Harmful prompts, jailbreaks, poisoned RAG documents.", "Azure pitch", "Foundry/OpenAI guardrails apply content filtering to prompts and completions; Prompt Shields detect user prompt attacks and document attacks; groundedness detects responses not supported by source material.", "Proof", "Jailbreak and poisoned-doc prompts block before unsafe model use.", "Content Safety overview; Prompt Shields quickstart; Foundry guardrails docs", TEAL),
        ("M3 · PII Protection", "Azure AI Language", "Exploit", "Prompt or tool output contains SSN/card/account data and leaks into responses or logs.", "Azure pitch", "Azure AI Language PII detection identifies sensitive entities so the app can transform data before model calls, responses, and logs. This is a data pipeline control, not just a model block.", "Proof", "The same answer is returned with sensitive values redacted and audit-safe events.", "Azure AI Language PII detection overview", AMBER),
        ("M4 · Tools, MCP, HITL, Code", "least privilege and action safety", "Exploit", "IDOR, SQLi, transfer without approval, unsafe report code, untrusted MCP tool, forged handoff.", "Azure pitch", "Foundry agents expose scoped hosted tools; PostgreSQL/RLS and MCP allow-lists narrow what an agent can do; human approval gates state changes; Code Interpreter provides a sandboxed execution boundary.", "Proof", "Read-only scope, approval-required transfer, sandbox blocks imports, MCP state-changing probe is denied.", "Foundry agents/tools; Azure Database for PostgreSQL; Code Interpreter; MCP tool scoping", GREEN),
        ("M5 · Identity + Search ACLs", "Entra OBO and AI Search document security", "Exploit", "Editable customer ID/groups return another customer or restricted docs.", "Azure pitch", "OBO propagates the signed-in user to downstream APIs. AI Search trims retrieval with a filterable `group_ids` field and `search.in()` over caller groups, so the RAG layer cannot see unauthorized documents.", "Proof", "Learner sees retail/public docs only; manager sees manager/private-client docs.", "Microsoft identity platform OBO; Azure AI Search security filters", BLUE),
        ("M6 · Gateway + Secure Runtime", "APIM, observability, Defender", "Exploit", "No throttling, exposed model/tool endpoints, verbose runtime errors, weak audit trail.", "Azure pitch", "APIM centralizes model/tool access with token limit policy, auth, logging, and key isolation. Monitor and Defender add operational evidence and threat detection around the workload.", "Proof", "Burst requests hit token budget; safe errors replace stack/detail leakage; logs show governed calls.", "APIM Azure OpenAI token limit policy; Azure Monitor; Defender for Cloud", TEAL),
        ("M7-M11 · Governance + Assurance", "AGT, groundedness, evals, Purview, red team", "Exploit", "Controls drift, poisoned data persists, no regression gate, tenant data policies missing, attacks untested.", "Azure pitch", "Governance inventories agents/tools, evaluations measure quality and safety, Purview DSPM/DLP extends data controls across the tenant, and red teaming turns attacks into repeatable coverage.", "Proof", "Policy check/eval/red-team reports become the gate before a lab fix is called done.", "Purview AI protection; Azure AI Foundry evaluations; red teaming transparency", CORAL),
    ]
    for title, subtitle, h1, b1, h2, b2, h3, b3, src, accent in layer_slides:
        s = prs.slides.add_slide(blank)
        add_bg(s)
        add_title(s, title, subtitle)
        add_card(s, 0.75, 1.55, 3.65, 3.7, h1, b1, CORAL, 15, 11.2)
        add_card(s, 4.85, 1.55, 3.65, 3.7, h2, b2, accent, 15, 10.8)
        add_card(s, 8.95, 1.55, 3.65, 3.7, h3, b3, GREEN, 15, 11.2)
        add_flow(s, ["run exploit", "enable layer", "rerun exploit", "capture proof"], 1.15, 5.75, 11.0, [CORAL, accent, BLUE, GREEN])
        source_footer(s, src)
        add_footer(s)

    # 12
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Walkthrough Agenda", "efficient presenter flow")
    items = [
        ("0-5 min", "Objective + app boundary", "Use slides 1-3. Make the finance/PII risk concrete."),
        ("5-20 min", "Baseline tour", "Show V1-V11 map and run 3 fast exploits: jailbreak, PII, IDOR."),
        ("20-55 min", "Core Azure layers", "Foundry guardrails, PII, tool/MCP/HITL, Entra/Search, APIM/runtime."),
        ("55-70 min", "Assurance", "Governance, evaluations, Purview, red-team as the regression story."),
        ("70+ min", "Hands-on lab", "Switch to MOAW and have learners run exploit -> fix -> verify."),
    ]
    for i, (time, head, body) in enumerate(items):
        y = 1.35 + i * 1.05
        add_pill(s, time, 0.75, y, 1.25, [TEAL, BLUE, AMBER, CORAL, GREEN][i])
        add_text(s, head, 2.2, y + 0.02, 3.5, 0.3, 14, INK, True, margin=False)
        add_text(s, body, 5.6, y + 0.02, 6.7, 0.38, 11, MUTED, margin=False)
    add_footer(s)

    # 13
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK)
    add_title(s, "Boundary Placement", "closing", dark=True)
    add_text(s, "Security for agentic apps is boundary placement.", 0.78, 1.6, 8.0, 0.55, 24, WHITE, True, TITLE_FONT, margin=False)
    add_card(s, 0.85, 2.65, 3.6, 2.1, "Platform first", "Use Foundry, Entra, Search ACLs, APIM, Purview, Monitor, and Defender where they can enforce controls outside app code.", TEAL, 15, 11)
    add_card(s, 4.85, 2.65, 3.6, 2.1, "App layer where needed", "Redact PII, inspect tool/MCP output, validate inter-agent handoffs, and gate state-changing actions.", AMBER, 15, 11)
    add_card(s, 8.85, 2.65, 3.6, 2.1, "Proof beats posture", "Every module ends with before/after evidence: blocked prompt, trimmed docs, denied tool, rate limit, eval or red-team result.", GREEN, 15, 11)
    add_text(s, "MOAW: docs/workshop.md", 0.85, 6.42, 4.0, 0.24, 9, RGBColor(199, 217, 213), margin=False)
    add_text(s, "Hosted: ca-app-zava8zn5p.yellowrock-178fa456.swedencentral.azurecontainerapps.io", 0.85, 6.72, 8.6, 0.24, 8.5, RGBColor(199, 217, 213), margin=False)

    # 14
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Microsoft Learn Sources", "for speaker backup")
    refs = [
        "Azure AI Content Safety: content filters, Prompt Shields, groundedness, protected material",
        "Azure AI Language: PII detection and redaction pipeline",
        "Azure AI Search: document-level security pattern with group_ids and search.in()",
        "Microsoft identity platform: OAuth 2.0 On-Behalf-Of flow",
        "Azure API Management: Azure OpenAI token limit and gateway policies",
        "Microsoft Purview: protect and govern data for generative AI apps",
        "Azure AI Foundry evaluations and red-team safety evaluation guidance",
    ]
    for i, ref in enumerate(refs):
        y = 1.4 + i * 0.68
        add_pill(s, str(i + 1), 0.8, y, 0.42, [TEAL, AMBER, BLUE, GREEN, CORAL, TEAL, BLUE][i])
        add_text(s, ref, 1.35, y + 0.02, 10.8, 0.28, 12.2, INK, margin=False)
    add_footer(s, "Keep detailed docs in MOAW; keep the PPT as the talk track.")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    build_deck()
